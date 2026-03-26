#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Flappy Bird Pygame - PowerPoint Esityksen Luominen
Luo professionaalisen PowerPoint-esityksen projektista
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """Luo otsikkoslaidin"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Tausta väri
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 71, 42)  # Tummanvihreä
    
    # Otsikko
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Alaotsikko
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(32)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(61, 106, 77)  # Vaaleanvihreä
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def create_content_slide(prs, title, content_list, background_color=RGBColor(245, 245, 245)):
    """Luo sisältöslaidin"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Tausta
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = background_color
    
    # Otsikko palkki
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(26, 71, 42)  # Tummanvihreä
    title_shape.line.color.rgb = RGBColor(26, 71, 42)
    
    # Otsikko teksti
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title_frame.margin_left = Inches(0.3)
    
    # Sisältö
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.3))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(12)
        p.space_after = Pt(12)
        p.level = 0
    
    return slide

def create_powerpoint_presentation():
    """Luo koko PowerPoint esityksen"""
    
    # Luo esitys
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # SLIDE 1: Otsikkoslaide
    create_title_slide(prs, 
                      "FLAPPY BIRD PYGAME",
                      "Pelikehityksen Matka - Python & Pygame")
    
    # SLIDE 2: Mikä on Flappy Bird?
    create_content_slide(prs,
                        "📱 Mikä on Flappy Bird?",
                        [
                            "✓ Klassinen 2D-peli jossa lintu lentää",
                            "✓ Yksinkertainen mutta haastava mekaniiikka",
                            "✓ Addiktiivinen gameplay - helppo oppia, vaikea hallita",
                            "✓ Täydellinen oppimisprojekti pelikehitykseen",
                            "✓ Koostuu: fysiikka, grafiikka, äänet, UI"
                        ])
    
    # SLIDE 3: Tekniikka
    create_content_slide(prs,
                        "🛠️ Tekniikka & Työkalut",
                        [
                            "• Ohjelmointikieli: Python 3.13",
                            "• Pelimoottori: Pygame 2.5.2",
                            "• Matematiikka: NumPy (äänet)",
                            "• Arkkitehtuuri: Tilakoneen malli",
                            "• Koodi: 451 rivoa + apumoduulit"
                        ])
    
    # SLIDE 4: Kehitysmatka
    create_content_slide(prs,
                        "📅 6 Viikon Kehitysmatka",
                        [
                            "Viikko 1: 🏗️ Perustukset & Arkkitehtuuri",
                            "Viikko 2: 🎨 Omaisuuksien generointi (ääni + grafiikka)",
                            "Viikko 3: ⚙️ Fysiikka & Liike",
                            "Viikko 4: 🖼️ Käyttöliittymä",
                            "Viikko 5: ⚡ Optimointi & Virheenkäsittely",
                            "Viikko 6: ✅ Testaus & Viimeistely"
                        ])
    
    # SLIDE 5: Viikko 1 - Perustukset
    create_content_slide(prs,
                        "🏗️ Viikko 1: Perustukset",
                        [
                            "✓ Virtuaaliympäristön asennus (.venv)",
                            "✓ Pygame & NumPy paketit asennettu",
                            "✓ Tilakoneen rakentaminen (3 tilaa)",
                            "✓ Bird ja Pipe -luokkien suunnittelu",
                            "✓ Pääsilmukan rakentaminen",
                            "⏱️ Aika: 3-4 tuntia"
                        ])
    
    # SLIDE 6: Viikko 2 - Omaisuudet
    create_content_slide(prs,
                        "🎨 Viikko 2: Omaisuuksien Generointi",
                        [
                            "✓ Äänitehosteet NumPy:n avulla generoidut",
                            "✓ Wing ääni (900 Hz hyppäämiseen)",
                            "✓ Swoosh ääni (taajuus sweep putken ohitykseen)",
                            "✓ Point ääni (1400 Hz pisteeseen)",
                            "✓ Grafiikat Pygame:llä piirretty (lintu, putket, tausta)",
                            "⏱️ Aika: 4-5 tuntia"
                        ])
    
    # SLIDE 7: Viikko 3 - Fysiikka
    create_content_slide(prs,
                        "⚙️ Viikko 3: Fysiikka & Törmäykset",
                        [
                            "✓ Painovoima toteutettu (0.4 px/frame²)",
                            "✓ Hyppääminen (Jump Strength = -6)",
                            "✓ Putken liike vasemmalle (-2 px/frame)",
                            "✓ Törmäystarkistus (putki, maa, väli)",
                            "✓ Pisteen laskeminen automaattisesti",
                            "⏱️ Aika: 5-6 tuntia (säätäminen vei aikaa!)"
                        ])
    
    # SLIDE 8: Viikko 4 - UI
    create_content_slide(prs,
                        "🖼️ Viikko 4: Käyttöliittymä",
                        [
                            "✓ Päävalikko - animoitu lintu, ohjeet",
                            "✓ Pelaamisen näyttö - putket, lintu, pisteet",
                            "✓ Game Over näyttö - lopputulokset",
                            "✓ Lintun rotaatio nopeuden mukaisesti",
                            "✓ Näytöiden sujuva vaihtaminen",
                            "⏱️ Aika: 3-4 tuntia"
                        ])
    
    # SLIDE 9: Viikko 5 - Optimointi
    create_content_slide(prs,
                        "⚡ Viikko 5: Optimointi & Kestävyys",
                        [
                            "✓ DummySound -luokka virheenkäsittelyyn",
                            "✓ Try/except blokit äänien lataamiselle",
                            "✓ Turhan piirtämisen poistaminen (-70% CPU)",
                            "✓ Objektien muistinhallinta (-50% RAM)",
                            "✓ FPS vakaa 60 ruutua/sekunti",
                            "⏱️ Aika: 2-3 tuntia"
                        ])
    
    # SLIDE 10: Viikko 6 - Testaus
    create_content_slide(prs,
                        "✅ Viikko 6: Testaus & Dokumentaatio",
                        [
                            "✓ 6 testisarjaa - kaikki läpäisty",
                            "✓ Pelin käynnistys, fysiikka, törmäykset",
                            "✓ Pisteet, uudelleenyritys, käyttöliittymä",
                            "✓ Tekninen raportti kirjoitettu (~1500 rivoa)",
                            "✓ Esitysohjeistus valmistettu",
                            "⏱️ Aika: 4-5 tuntia"
                        ])
    
    # SLIDE 11: Arkkitehtuuri
    create_content_slide(prs,
                        "🔧 Arkkitehtuuri: Tilakoneen Malli",
                        [
                            "TILA 0: Päävalikko (animoitu lintu)",
                            "↓ (Välilyönti)",
                            "TILA 1: Peli Käynnissä (fysiikka + törmäykset)",
                            "↓ (Osuma)",
                            "TILA 2: Game Over (tulokset)",
                            "↓ (Välilyönti)",
                            "↻ Takaisin TILAAN 0"
                        ])
    
    # SLIDE 12: Haasteet & Ratkaisut
    create_content_slide(prs,
                        "💡 Haasteet & Ratkaisut",
                        [
                            "Haaste: Fysiikan säätäminen",
                            "→ Ratkaisu: Iteratiivinen testaus ja tuning",
                            "",
                            "Haaste: Ääni-järjestelmän epävarmuus",
                            "→ Ratkaisu: DummySound fallback-mekanismi",
                            "",
                            "Haaste: Suorituskyvyn ongelmat",
                            "→ Ratkaisu: Optimointi ja profilointi"
                        ])
    
    # SLIDE 13: Saavutukset
    create_content_slide(prs,
                        "🏆 Saavutukset",
                        [
                            "✅ Täysin toimiva peli - 451 rivoa koodia",
                            "✅ Ammattimainen arkkitehtuuri",
                            "✅ Proseduraalinen äänien generointi",
                            "✅ 60 FPS vakaa suorituskyky",
                            "✅ Perusteellinen dokumentaatio",
                            "✅ Testaus suoritettu - laatu varmistettu"
                        ])
    
    # SLIDE 14: Opitut Asiat
    create_content_slide(prs,
                        "📚 Opitut Asiat",
                        [
                            "• Pelimoottoreiden perusteet",
                            "• Fysiikan toteutus pelissä",
                            "• Olio-ohjelmointi käytännössä",
                            "• Käyttöliittymän suunnittelu",
                            "• Virheenkäsittely ja robustisuus",
                            "• Testaus ja laadunvarmistus"
                        ])
    
    # SLIDE 15: Tulevaisuus
    create_content_slide(prs,
                        "🚀 Tulevaisuuden Kehitys",
                        [
                            "💡 Lisää vaikeustasoja (Easy, Medium, Hard)",
                            "💡 Parhaat pisteet -taulukko (high scores)",
                            "💡 Power-upeja (nopeus, suoja, magneetit)",
                            "💡 Pelimuodot (Time Attack, Survival)",
                            "💡 Verkkopelaajille (multiplayer)",
                            "💡 Mobiili-versio (Android/iOS)"
                        ])
    
    # SLIDE 16: Yhteenveto
    create_content_slide(prs,
                        "📝 Yhteenveto",
                        [
                            "🎮 Täysin toimiva pelisovellu",
                            "📊 6 viikon intensiivinen kehitys",
                            "💻 451 rivoa Python-koodia",
                            "🧪 Testattu ja optimoitu",
                            "📚 Perusteellisesti dokumentoitu",
                            "🎓 Paljon opittavaa pelikehityksestä"
                        ])
    
    # SLIDE 17: Kiitos
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 71, 42)
    
    kiitos_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    kiitos_frame = kiitos_box.text_frame
    kiitos_frame.text = "KIITOS!"
    kiitos_frame.paragraphs[0].font.size = Pt(88)
    kiitos_frame.paragraphs[0].font.bold = True
    kiitos_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    kiitos_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    kysymykset_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1.5))
    kysymykset_frame = kysymykset_box.text_frame
    kysymykset_frame.text = "Kysymyksiä?"
    kysymykset_frame.paragraphs[0].font.size = Pt(44)
    kysymykset_frame.paragraphs[0].font.color.rgb = RGBColor(61, 106, 77)
    kysymykset_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Tallenna
    output_path = "FLAPPY_BIRD_ESITYS.pptx"
    prs.save(output_path)
    
    print(f"✅ PowerPoint-esitys luotu onnistuneesti!")
    print(f"📄 Tiedosto: {output_path}")
    print(f"📊 Slaideja yhteensä: {len(prs.slides)}")
    print(f"🎨 Teema: Vihreä/Teal (ammattimainen)")
    print(f"✓ Esitys on valmis käyttöön!")

if __name__ == "__main__":
    create_powerpoint_presentation()
