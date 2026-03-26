#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Flappy Bird Pygame - PowerPoint Esityksen Luominen KUVILLA
Luo professionaalisen PowerPoint-esityksen projektista + visualisointikuvat
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageDraw, ImageFont
import os

def create_architecture_image():
    """Luo tilakoneen arkkitehtuurikaavio"""
    img = Image.new('RGB', (800, 600), color=(245, 245, 245))
    draw = ImageDraw.Draw(img)
    
    # Värit
    dark_green = (26, 71, 42)
    light_green = (61, 106, 77)
    white = (255, 255, 255)
    black = (0, 0, 0)
    
    # Piirä tilat
    # Tila 0 - Päävalikko
    draw.rectangle([50, 100, 200, 200], fill=dark_green, outline=black, width=3)
    draw.text((70, 140), "TILA 0\nPÄÄVALIKKO", fill=white, font=None)
    
    # Tila 1 - Peli käynnissä
    draw.rectangle([350, 100, 500, 200], fill=light_green, outline=black, width=3)
    draw.text((360, 140), "TILA 1\nPELI KÄYNNISSÄ", fill=white, font=None)
    
    # Tila 2 - Game Over
    draw.rectangle([600, 100, 750, 200], fill=dark_green, outline=black, width=3)
    draw.text((620, 140), "TILA 2\nGAME OVER", fill=white, font=None)
    
    # Piirä nuolet
    # Tila 0 → Tila 1
    draw.line([(200, 150), (350, 150)], fill=black, width=3)
    draw.text((260, 120), "VÄLILYÖNTI", fill=black, font=None)
    
    # Tila 1 → Tila 2
    draw.line([(500, 150), (600, 150)], fill=black, width=3)
    draw.text((540, 120), "TÖRMÄYS", fill=black, font=None)
    
    # Tila 2 → Tila 0
    draw.arc([(550, 250), (750, 450)], 0, 180, fill=black, width=3)
    draw.text((600, 380), "VÄLILYÖNTI", fill=black, font=None)
    
    img.save('img_architecture.png')
    return 'img_architecture.png'

def create_physics_chart():
    """Luo fysiikan visualisointikaavio"""
    img = Image.new('RGB', (800, 600), color=(245, 245, 245))
    draw = ImageDraw.Draw(img)
    
    # Akselit
    draw.line([(50, 550), (750, 550)], fill=(0, 0, 0), width=2)  # X-akseli
    draw.line([(50, 50), (50, 550)], fill=(0, 0, 0), width=2)    # Y-akseli
    
    # Akselin tekstit
    draw.text((750, 530), "Aika", fill=(0, 0, 0))
    draw.text((30, 20), "Y-sijainti", fill=(0, 0, 0))
    
    # Piirä paraabeli (painovoiman vaikutus)
    points = []
    for x in range(100, 700, 5):
        # Yksinkertainen paraabeli: y = (x-400)²/200
        y = int(550 - ((x - 400) ** 2) / 200)
        points.append((x, y))
    
    # Piirä käyrä
    for i in range(len(points) - 1):
        draw.line([points[i], points[i+1]], fill=(26, 71, 42), width=3)
    
    # Merkitse hyppyä
    draw.text((350, 300), "Hyppää\n↑", fill=(0, 0, 200))
    draw.text((500, 500), "Putoaa\n↓", fill=(200, 0, 0))
    
    img.save('img_physics.png')
    return 'img_physics.png'

def create_game_screenshot():
    """Luo pelinäytön simulaatio"""
    img = Image.new('RGB', (800, 600), color=(135, 206, 235))  # Sininen taivas
    draw = ImageDraw.Draw(img)
    
    # Maa
    draw.rectangle([0, 550, 800, 600], fill=(34, 139, 34))
    
    # Lintu (keltainen ympyrä)
    draw.ellipse([300, 250, 350, 300], fill=(255, 255, 0), outline=(0, 0, 0), width=2)
    draw.ellipse([320, 270, 325, 275], fill=(0, 0, 0))  # Silmä
    
    # Putket (vihreät)
    # Ylempi putki
    draw.rectangle([600, 0, 650, 250], fill=(34, 139, 34), outline=(0, 0, 0), width=2)
    # Alempi putki
    draw.rectangle([600, 350, 650, 600], fill=(34, 139, 34), outline=(0, 0, 0), width=2)
    
    # Putket (toiset)
    draw.rectangle([400, 0, 450, 200], fill=(34, 139, 34), outline=(0, 0, 0), width=2)
    draw.rectangle([400, 300, 450, 600], fill=(34, 139, 34), outline=(0, 0, 0), width=2)
    
    # Pisteet näytön yläosassa
    draw.text((50, 20), "Pisteet: 5", fill=(0, 0, 0))
    
    img.save('img_gameplay.png')
    return 'img_gameplay.png'

def create_weekly_summary():
    """Luo viikkoittaisen yhteenvedon"""
    img = Image.new('RGB', (800, 600), color=(245, 245, 245))
    draw = ImageDraw.Draw(img)
    
    weeks = [
        ("Viikko 1", "Perustukset"),
        ("Viikko 2", "Omaisuudet"),
        ("Viikko 3", "Fysiikka"),
        ("Viikko 4", "UI"),
        ("Viikko 5", "Optimointi"),
        ("Viikko 6", "Testaus"),
    ]
    
    y_pos = 80
    box_width = 120
    box_height = 80
    spacing = 20
    
    colors = [
        (26, 71, 42),
        (61, 106, 77),
        (95, 141, 111),
        (26, 71, 42),
        (61, 106, 77),
        (95, 141, 111),
    ]
    
    x_pos = 80
    for i, (week, task) in enumerate(weeks):
        color = colors[i]
        # Piirä boksi
        draw.rectangle([x_pos, y_pos, x_pos + box_width, y_pos + box_height],
                      fill=color, outline=(0, 0, 0), width=2)
        # Teksti
        draw.text((x_pos + 10, y_pos + 15), week, fill=(255, 255, 255))
        draw.text((x_pos + 5, y_pos + 40), task, fill=(255, 255, 255))
        
        x_pos += box_width + spacing
    
    img.save('img_weekly_summary.png')
    return 'img_weekly_summary.png'

def create_code_example():
    """Luo koodiesimerkin"""
    img = Image.new('RGB', (800, 600), color=(30, 30, 30))  # Tumma tausta
    draw = ImageDraw.Draw(img)
    
    code_lines = [
        "# Tilakoneen perusta",
        "game_state = 0",
        "",
        "while running:",
        "    if game_state == 0:",
        "        title_screen()",
        "    elif game_state == 1:",
        "        game_running()",
        "    elif game_state == 2:",
        "        game_over_screen()",
    ]
    
    y_pos = 50
    for line in code_lines:
        if line.startswith("#"):
            color = (100, 200, 100)  # Vihreä kommentille
        elif "game_state" in line:
            color = (200, 100, 255)  # Violetti muuttujille
        else:
            color = (255, 255, 255)  # Valkoinen muulle
        
        draw.text((50, y_pos), line, fill=color)
        y_pos += 40
    
    img.save('img_code_example.png')
    return 'img_code_example.png'

def create_performance_chart():
    """Luo suorituskyvyn kaavio"""
    img = Image.new('RGB', (800, 600), color=(245, 245, 245))
    draw = ImageDraw.Draw(img)
    
    # Otsikko
    draw.text((300, 30), "Optimointi - Ennen vs Jälkeen", fill=(0, 0, 0))
    
    metrics = [
        ("CPU %", 12, 3),
        ("RAM MB", 60, 30),
        ("FPS", 45, 60),
    ]
    
    bar_height = 30
    bar_width = 60
    x_start = 100
    y_start = 150
    
    for i, (metric, before, after) in enumerate(metrics):
        y_pos = y_start + i * 120
        
        # Metriikan nimi
        draw.text((x_start - 80, y_pos), metric, fill=(0, 0, 0))
        
        # Ennen (punainen)
        draw.rectangle([x_start, y_pos, x_start + before * 3, y_pos + bar_height],
                      fill=(255, 100, 100), outline=(0, 0, 0), width=2)
        draw.text((x_start + before * 3 + 10, y_pos), f"Ennen: {before}", fill=(0, 0, 0))
        
        # Jälkeen (vihreä)
        y_after = y_pos + bar_height + 10
        draw.rectangle([x_start, y_after, x_start + after * 3, y_after + bar_height],
                      fill=(100, 255, 100), outline=(0, 0, 0), width=2)
        draw.text((x_start + after * 3 + 10, y_after), f"Jälkeen: {after}", fill=(0, 0, 0))
    
    img.save('img_performance.png')
    return 'img_performance.png'

def add_title_slide(prs, title, subtitle):
    """Luo otsikkoslaidin"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 71, 42)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(32)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(61, 106, 77)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide_with_image(prs, title, content_list, image_path=None, image_position="right"):
    """Luo sisältöslaidin kuvalla"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # Otsikko palkki
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(26, 71, 42)
    title_shape.line.color.rgb = RGBColor(26, 71, 42)
    
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title_frame.margin_left = Inches(0.3)
    
    # Sisältö
    if image_position == "right":
        content_width = 4.5
        content_left = 0.5
        image_left = 5.3
    else:
        content_width = 4.5
        content_left = 4.2
        image_left = 0.5
    
    content_box = slide.shapes.add_textbox(Inches(content_left), Inches(1.2), 
                                          Inches(content_width), Inches(5.3))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    # Lisää kuva
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(image_left), Inches(1.2), 
                                width=Inches(4), height=Inches(5.5))

def create_powerpoint_with_images():
    """Luo PowerPoint esityksen kuvilla"""
    
    print("📸 Luodaan visualisointikuvia...")
    
    # Luo kuvat
    arch_img = create_architecture_image()
    physics_img = create_physics_chart()
    gameplay_img = create_game_screenshot()
    weekly_img = create_weekly_summary()
    code_img = create_code_example()
    perf_img = create_performance_chart()
    
    print("✅ Kuvat luotu!")
    print("🎬 Luodaan PowerPoint-esitystä...")
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # SLIDE 1: Otsikko
    add_title_slide(prs, "FLAPPY BIRD PYGAME", 
                   "Pelikehityksen Matka - Python & Pygame")
    
    # SLIDE 2: Mikä on Flappy Bird
    add_content_slide_with_image(prs,
        "📱 Mikä on Flappy Bird?",
        [
            "✓ Klassinen 2D-peli",
            "✓ Yksinkertainen mutta haastava",
            "✓ Täydellinen oppimisprojekti",
            "✓ Koostuu: fysiikka, grafiikka, äänet, UI"
        ],
        gameplay_img,
        "right")
    
    # SLIDE 3: Tekniikka
    add_content_slide_with_image(prs,
        "🛠️ Tekniikka & Työkalut",
        [
            "• Python 3.13",
            "• Pygame 2.5.2",
            "• NumPy (äänet)",
            "• Tilakoneen arkkitehtuuri",
            "• 451 rivoa koodia"
        ],
        code_img,
        "right")
    
    # SLIDE 4: Kehitysmatka
    add_content_slide_with_image(prs,
        "📅 6 Viikon Kehitysmatka",
        [
            "Viikko 1: 🏗️ Perustukset",
            "Viikko 2: 🎨 Omaisuudet",
            "Viikko 3: ⚙️ Fysiikka",
            "Viikko 4: 🖼️ Käyttöliittymä",
            "Viikko 5: ⚡ Optimointi",
            "Viikko 6: ✅ Testaus"
        ],
        weekly_img,
        "right")
    
    # SLIDE 5: Arkkitehtuuri
    add_content_slide_with_image(prs,
        "🔧 Arkkitehtuuri: Tilakoneen Malli",
        [
            "TILA 0: Päävalikko",
            "↓ (Välilyönti)",
            "TILA 1: Peli käynnissä",
            "↓ (Törmäys)",
            "TILA 2: Game Over",
            "↓ (Välilyönti)",
            "↻ Takaisin"
        ],
        arch_img,
        "right")
    
    # SLIDE 6: Fysiikka
    add_content_slide_with_image(prs,
        "⚙️ Fysiikka & Liike",
        [
            "✓ Painovoima: 0.4 px/frame²",
            "✓ Hyppääminen: -6 px/frame",
            "✓ Putken nopeus: -2 px/frame",
            "✓ 60 FPS vakaa",
            "✓ Säätäminen: iteratiivinen testaus"
        ],
        physics_img,
        "right")
    
    # SLIDE 7: Haasteet & Ratkaisut
    add_content_slide_with_image(prs,
        "💡 Haasteet & Ratkaisut",
        [
            "Fysiikan säätäminen",
            "→ Iteratiivinen testaus",
            "",
            "Ääni-järjestelmä",
            "→ DummySound fallback",
            "",
            "Suorituskyky",
            "→ Optimointi"
        ],
        perf_img,
        "right")
    
    # SLIDE 8: Saavutukset
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = RGBColor(26, 71, 42)
    title_shape.line.color.rgb = RGBColor(26, 71, 42)
    
    title_frame = title_shape.text_frame
    title_frame.text = "🏆 Saavutukset & Tulokset"
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    title_frame.margin_left = Inches(0.3)
    
    achievements = [
        "✅ Täysin toimiva peli",
        "✅ 451 rivoa koodia",
        "✅ Proseduraalinen äänet",
        "✅ 60 FPS vakaa",
        "✅ -70% CPU käyttö",
        "✅ -50% RAM käyttö",
        "✅ 8/8 testia läpäisty",
        "✅ Dokumentaatio valmis"
    ]
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    text_frame = content_box.text_frame
    
    for i, achievement in enumerate(achievements):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = achievement
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(26, 71, 42)
        p.font.bold = True
        p.space_before = Pt(10)
        p.space_after = Pt(10)
    
    # SLIDE 9: Kiitos
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 71, 42)
    
    kiitos_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    kiitos_frame = kiitos_box.text_frame
    kiitos_frame.text = "KIITOS!"
    kiitos_frame.paragraphs[0].font.size = Pt(88)
    kiitos_frame.paragraphs[0].font.bold = True
    kiitos_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    kiitos_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    kysymykset_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1.5))
    kysymykset_frame = kysymykset_box.text_frame
    kysymykset_frame.text = "Kysymyksiä?"
    kysymykset_frame.paragraphs[0].font.size = Pt(44)
    kysymykset_frame.paragraphs[0].font.color.rgb = RGBColor(61, 106, 77)
    kysymykset_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Tallenna
    output_path = "FLAPPY_BIRD_ESITYS.pptx"
    prs.save(output_path)
    
    print(f"\n✅ PowerPoint-esitys luotu onnistuneesti!")
    print(f"📄 Tiedosto: {output_path}")
    print(f"📊 Slaideja yhteensä: {len(prs.slides)}")
    print(f"🎨 Teema: Vihreä/Teal (ammattimainen)")
    print(f"📸 Visualisointikuvat: 6 kuvaa")
    print(f"✓ Esitys on valmis käyttöön!")
    
    # Näytä kuvat
    print(f"\n📁 Luodut kuvat:")
    print(f"   1. {arch_img} - Tilakoneen arkkitehtuuri")
    print(f"   2. {physics_img} - Fysiikan visualisointi")
    print(f"   3. {gameplay_img} - Pelinäyttö")
    print(f"   4. {weekly_img} - Viikkoittainen yhteenveto")
    print(f"   5. {code_img} - Koodiesimerkit")
    print(f"   6. {perf_img} - Suorituskyvyn kaavio")

if __name__ == "__main__":
    create_powerpoint_with_images()
